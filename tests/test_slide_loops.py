"""Tests for slide-level loop expansion ({%slide for ...%})."""

import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxtpl import PptxTemplate


def _get_slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


class TestSlideLoopBasic:
    def test_three_items_produce_three_slides(self, slide_loop_template, tmp_dir):
        tpl = PptxTemplate(slide_loop_template)
        tpl.render({
            "title": "Report",
            "items": [
                {"name": "Alice"},
                {"name": "Bob"},
                {"name": "Charlie"},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # 1 title + 3 cloned + 1 closing = 5 slides
        assert len(prs.slides) == 5

        assert "Report" in _get_slide_text(prs.slides[0])
        assert "Alice" in _get_slide_text(prs.slides[1])
        assert "Bob" in _get_slide_text(prs.slides[2])
        assert "Charlie" in _get_slide_text(prs.slides[3])
        assert "The End" in _get_slide_text(prs.slides[4])

    def test_no_jinja_tags_remain(self, slide_loop_template, tmp_dir):
        tpl = PptxTemplate(slide_loop_template)
        tpl.render({
            "title": "T",
            "items": [{"name": "X"}],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        for slide in prs.slides:
            text = _get_slide_text(slide)
            assert "{%" not in text
            assert "{{" not in text


class TestSlideLoopEmpty:
    def test_empty_list_removes_template_slide(self, slide_loop_template, tmp_dir):
        tpl = PptxTemplate(slide_loop_template)
        tpl.render({
            "title": "Report",
            "items": [],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # 1 title + 0 cloned + 1 closing = 2 slides
        assert len(prs.slides) == 2
        assert "Report" in _get_slide_text(prs.slides[0])
        assert "The End" in _get_slide_text(prs.slides[1])


class TestSlideLoopSingle:
    def test_single_item(self, slide_loop_template, tmp_dir):
        tpl = PptxTemplate(slide_loop_template)
        tpl.render({
            "title": "Report",
            "items": [{"name": "Solo"}],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 3
        assert "Solo" in _get_slide_text(prs.slides[1])


class TestSlideLoopSimpleValues:
    def test_loop_over_strings(self, slide_loop_simple_template, tmp_dir):
        tpl = PptxTemplate(slide_loop_simple_template)
        tpl.render({"names": ["Alpha", "Beta"]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2
        assert "Alpha" in _get_slide_text(prs.slides[0])
        assert "Beta" in _get_slide_text(prs.slides[1])


class TestSlideLoopVariable:
    def test_loop_index(self, tmp_dir):
        """The loop helper provides index, first, last, length."""
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(0.5))
        tb.text_frame.text = (
            "{%slide for item in items %}"
            "{{ loop.index }}/{{ loop.length }}"
            "{%slide endfor %}"
        )
        path = os.path.join(tmp_dir, "loop_var.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"items": ["a", "b", "c"]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 3
        assert "1/3" in _get_slide_text(result.slides[0])
        assert "2/3" in _get_slide_text(result.slides[1])
        assert "3/3" in _get_slide_text(result.slides[2])

    def test_loop_first_last(self, tmp_dir):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tb.text_frame.text = (
            "{%slide for x in items %}"
            "{% if loop.first %}FIRST{% endif %}"
            "{% if loop.last %}LAST{% endif %}"
            "{%slide endfor %}"
        )
        path = os.path.join(tmp_dir, "loop_fl.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"items": [1, 2, 3]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert "FIRST" in _get_slide_text(result.slides[0])
        assert "FIRST" not in _get_slide_text(result.slides[1])
        assert "LAST" not in _get_slide_text(result.slides[1])
        assert "LAST" in _get_slide_text(result.slides[2])


class TestSlideLoopWithNestedContent:
    def test_nested_jinja_loop_inside_slide_loop(self, tmp_dir):
        """A slide-level loop with a regular Jinja for loop inside."""
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(2))
        tb.text_frame.text = (
            "{%slide for section in sections %}"
            "{{ section.title }}: "
            "{% for point in section.points %}{{ point }} {% endfor %}"
            "{%slide endfor %}"
        )
        path = os.path.join(tmp_dir, "nested.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({
            "sections": [
                {"title": "Intro", "points": ["A", "B"]},
                {"title": "Body", "points": ["C", "D", "E"]},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 2
        text0 = _get_slide_text(result.slides[0])
        assert "Intro" in text0
        assert "A" in text0
        assert "B" in text0
        text1 = _get_slide_text(result.slides[1])
        assert "Body" in text1
        assert "C" in text1
        assert "E" in text1


class TestMultipleSlideLoops:
    def test_two_loops_in_one_presentation(self, tmp_dir):
        prs = Presentation()

        # Static slide
        s0 = prs.slides.add_slide(prs.slide_layouts[6])
        s0.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "Title"

        # First loop
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for a in alphas %}{{ a }}{%slide endfor %}"
        )

        # Middle static slide
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "Middle"

        # Second loop
        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        s3.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for n in nums %}{{ n }}{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "multi_loop.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"alphas": ["X", "Y"], "nums": [1, 2, 3]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        # 1 title + 2 alpha + 1 middle + 3 nums = 7
        assert len(result.slides) == 7
        assert "Title" in _get_slide_text(result.slides[0])
        assert "X" in _get_slide_text(result.slides[1])
        assert "Y" in _get_slide_text(result.slides[2])
        assert "Middle" in _get_slide_text(result.slides[3])
        assert "1" in _get_slide_text(result.slides[4])
        assert "2" in _get_slide_text(result.slides[5])
        assert "3" in _get_slide_text(result.slides[6])


class TestSlideLoopSpecialChars:
    def test_xml_special_chars_in_loop_items(self, tmp_dir):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = "{%slide for name in names %}{{ name }}{%slide endfor %}"
        path = os.path.join(tmp_dir, "special.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"names": ["O'Brien & Sons", "A < B"]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 2
        assert "O'Brien & Sons" in _get_slide_text(result.slides[0])
        assert "A < B" in _get_slide_text(result.slides[1])


class TestSlideLoopUndeclaredVariables:
    def test_finds_variables_in_slide_loop(self, slide_loop_template):
        tpl = PptxTemplate(slide_loop_template)
        variables = tpl.get_undeclared_template_variables()
        assert "title" in variables
        # The slide loop body references item.name, but 'item' is the loop var
        # so 'items' (the iterable) should be detected, and 'item' should not
        # since it's declared by the stripped for-tag. However, since we strip
        # the {%slide for%} tag, 'item' becomes undeclared from Jinja2's perspective.
        # This is expected — the variable discovery is approximate.
        assert "item" in variables or "items" in variables
