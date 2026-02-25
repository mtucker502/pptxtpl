"""Tests for slide-level conditional inclusion ({%slide if ...%})."""

import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxtpl import PptxTemplate


def _get_slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


class TestConditionTrue:
    def test_slide_kept_when_condition_truthy(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "Report",
            "financials": {"revenue": "$4.2M"},
            "feedback": "Great quarter!",
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # All 4 slides kept
        assert len(prs.slides) == 4
        assert "Report" in _get_slide_text(prs.slides[0])
        assert "$4.2M" in _get_slide_text(prs.slides[1])
        assert "Great quarter!" in _get_slide_text(prs.slides[2])
        assert "The End" in _get_slide_text(prs.slides[3])

    def test_tags_stripped_from_kept_slide(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "T",
            "financials": {"revenue": "X"},
            "feedback": "Y",
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        for slide in prs.slides:
            text = _get_slide_text(slide)
            assert "{%slide" not in text
            assert "slide if" not in text
            assert "slide endif" not in text


class TestConditionFalse:
    def test_slide_removed_when_condition_falsy(self, conditional_slide_template, tmp_dir):
        """Both conditional slides removed when keys are missing/falsy."""
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({"title": "Report"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # Only title + closing = 2 slides
        assert len(prs.slides) == 2
        assert "Report" in _get_slide_text(prs.slides[0])
        assert "The End" in _get_slide_text(prs.slides[1])

    def test_none_value_removes_slide(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "Report",
            "financials": None,
            "feedback": None,
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2

    def test_empty_string_removes_slide(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "Report",
            "financials": "",
            "feedback": "",
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2


class TestMixedConditions:
    def test_some_true_some_false(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "Report",
            "financials": {"revenue": "$1M"},
            # feedback is missing → falsy
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # title + financials + closing = 3
        assert len(prs.slides) == 3
        assert "Report" in _get_slide_text(prs.slides[0])
        assert "$1M" in _get_slide_text(prs.slides[1])
        assert "The End" in _get_slide_text(prs.slides[2])

    def test_second_true_first_false(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "Report",
            # financials missing → removed
            "feedback": "Excellent",
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # title + feedback + closing = 3
        assert len(prs.slides) == 3
        assert "Report" in _get_slide_text(prs.slides[0])
        assert "Excellent" in _get_slide_text(prs.slides[1])
        assert "The End" in _get_slide_text(prs.slides[2])


class TestComplexExpression:
    def test_filter_expression(self, tmp_dir):
        """Condition using a Jinja2 filter: items|length > 0."""
        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "Title"

        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb.text_frame.text = (
            "{%slide if items|length > 0 %}"
            "Items: {% for i in items %}{{ i }} {% endfor %}"
            "{%slide endif %}"
        )
        path = os.path.join(tmp_dir, "complex_cond.pptx")
        prs.save(path)

        # Non-empty list → slide kept
        tpl = PptxTemplate(path)
        tpl.render({"items": ["a", "b"]})
        output = os.path.join(tmp_dir, "out_nonempty.pptx")
        tpl.save(output)
        result = Presentation(output)
        assert len(result.slides) == 2
        assert "a" in _get_slide_text(result.slides[1])

        # Empty list → slide removed
        tpl2 = PptxTemplate(path)
        tpl2.render({"items": []})
        output2 = os.path.join(tmp_dir, "out_empty.pptx")
        tpl2.save(output2)
        result2 = Presentation(output2)
        assert len(result2.slides) == 1
        assert "Title" in _get_slide_text(result2.slides[0])

    def test_boolean_expression(self, tmp_dir):
        """Condition with a boolean and/or expression."""
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = (
            "{%slide if show_a and show_b %}Both shown{%slide endif %}"
        )
        path = os.path.join(tmp_dir, "bool_cond.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"show_a": True, "show_b": True})
        output = os.path.join(tmp_dir, "out.pptx")
        tpl.save(output)
        result = Presentation(output)
        assert len(result.slides) == 1
        assert "Both shown" in _get_slide_text(result.slides[0])

        tpl2 = PptxTemplate(path)
        tpl2.render({"show_a": True, "show_b": False})
        output2 = os.path.join(tmp_dir, "out2.pptx")
        tpl2.save(output2)
        result2 = Presentation(output2)
        assert len(result2.slides) == 0


class TestCombinedWithSlideLoop:
    def test_conditional_gates_loop(self, tmp_dir):
        """A conditional slide followed by a loop slide — when condition is false,
        the conditional slide is removed and the loop still works."""
        prs = Presentation()

        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "Title"

        # Conditional slide
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = (
            "{%slide if show_intro %}Intro content{%slide endif %}"
        )

        # Loop slide
        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        s3.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for item in items %}{{ item }}{%slide endfor %}"
        )

        s4 = prs.slides.add_slide(prs.slide_layouts[6])
        s4.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "End"

        path = os.path.join(tmp_dir, "cond_loop.pptx")
        prs.save(path)

        # Condition false — intro removed, loop still expands
        tpl = PptxTemplate(path)
        tpl.render({"show_intro": False, "items": ["A", "B"]})
        output = os.path.join(tmp_dir, "out.pptx")
        tpl.save(output)
        result = Presentation(output)
        # Title + 2 loop slides + End = 4
        assert len(result.slides) == 4
        assert "Title" in _get_slide_text(result.slides[0])
        assert "A" in _get_slide_text(result.slides[1])
        assert "B" in _get_slide_text(result.slides[2])
        assert "End" in _get_slide_text(result.slides[3])

    def test_conditional_true_with_loop(self, tmp_dir):
        """Both conditional and loop present, condition true."""
        prs = Presentation()

        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide if show_intro %}Intro{%slide endif %}"
        )

        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for x in items %}{{ x }}{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "cond_loop2.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"show_intro": True, "items": ["Z"]})
        output = os.path.join(tmp_dir, "out.pptx")
        tpl.save(output)
        result = Presentation(output)
        # Intro + 1 loop slide = 2
        assert len(result.slides) == 2
        assert "Intro" in _get_slide_text(result.slides[0])
        assert "Z" in _get_slide_text(result.slides[1])

    def test_conditional_false_skips_loop_iterable(self, tmp_dir):
        """When condition is false, the loop's iterable doesn't need to exist."""
        prs = Presentation()

        # Conditional slide that wraps a loop reference
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = (
            "{%slide if show_details %}"
            "{% for d in details %}{{ d }}{% endfor %}"
            "{%slide endif %}"
        )

        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "End"

        path = os.path.join(tmp_dir, "cond_skip.pptx")
        prs.save(path)

        # details not in context, but slide is removed so no error
        tpl = PptxTemplate(path)
        tpl.render({"show_details": False})
        output = os.path.join(tmp_dir, "out.pptx")
        tpl.save(output)
        result = Presentation(output)
        assert len(result.slides) == 1
        assert "End" in _get_slide_text(result.slides[0])


class TestAllConditionsFalse:
    def test_only_static_slides_remain(self, conditional_slide_template, tmp_dir):
        tpl = PptxTemplate(conditional_slide_template)
        tpl.render({
            "title": "Only Title",
            "financials": False,
            "feedback": 0,
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2
        assert "Only Title" in _get_slide_text(prs.slides[0])
        assert "The End" in _get_slide_text(prs.slides[1])


class TestUndeclaredVariables:
    def test_discovers_variables_with_slide_if_tags(self, conditional_slide_template):
        tpl = PptxTemplate(conditional_slide_template)
        variables = tpl.get_undeclared_template_variables()
        assert "title" in variables
        # After stripping {%slide if%} tags, the body references remain
        assert "financials" in variables
        assert "feedback" in variables
