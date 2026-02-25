"""Fixtures that programmatically create .pptx templates for testing."""

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from lxml import etree


@pytest.fixture
def tmp_dir():
    """Provide a temporary directory that is cleaned up after the test."""
    with tempfile.TemporaryDirectory() as d:
        yield d


@pytest.fixture
def simple_template(tmp_dir):
    """Create a .pptx with a single slide containing {{ name }} placeholder."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.text = "Hello {{ name }}!"
    path = os.path.join(tmp_dir, "simple.pptx")
    prs.save(path)
    return path


@pytest.fixture
def multi_var_template(tmp_dir):
    """Create a .pptx with multiple template variables on one slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = txBox.text_frame
    tf.text = "{{ greeting }}, {{ name }}!"
    path = os.path.join(tmp_dir, "multi_var.pptx")
    prs.save(path)
    return path


@pytest.fixture
def conditional_template(tmp_dir):
    """Create a .pptx with a Jinja2 if/else block."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = txBox.text_frame
    tf.text = "{% if show_greeting %}Hello {{ name }}!{% endif %}"
    path = os.path.join(tmp_dir, "conditional.pptx")
    prs.save(path)
    return path


@pytest.fixture
def loop_template(tmp_dir):
    """Create a .pptx with a Jinja2 for loop."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = txBox.text_frame
    tf.text = "{% for item in items %}{{ item }} {% endfor %}"
    path = os.path.join(tmp_dir, "loop.pptx")
    prs.save(path)
    return path


@pytest.fixture
def multi_slide_template(tmp_dir):
    """Create a .pptx with template variables on multiple slides."""
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox1.text_frame.text = "Slide 1: {{ title }}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox2.text_frame.text = "Slide 2: {{ subtitle }}"

    path = os.path.join(tmp_dir, "multi_slide.pptx")
    prs.save(path)
    return path


@pytest.fixture
def fragmented_template(tmp_dir):
    """Create a .pptx where Jinja tags are manually split across multiple runs.

    This simulates PowerPoint's behavior of splitting text into separate
    <a:r> elements, which fragments Jinja2 delimiters.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]

    # Manually build fragmented runs: "Hello {" + "{" + " name " + "}" + "}!"
    # Clear the default run and build our own via XML manipulation
    p.text = ""
    p_element = p._element

    # Define the namespace
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    # Remove default empty run if present
    for r in p_element.findall("a:r", nsmap):
        p_element.remove(r)

    def add_run(parent, text):
        r = etree.SubElement(parent, f'{{{nsmap["a"]}}}r')
        t = etree.SubElement(r, f'{{{nsmap["a"]}}}t')
        t.text = text

    add_run(p_element, "Hello {")
    add_run(p_element, "{ name }")
    add_run(p_element, "}!")

    path = os.path.join(tmp_dir, "fragmented.pptx")
    prs.save(path)
    return path


@pytest.fixture
def table_template(tmp_dir):
    """Create a .pptx with a table containing template variables."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add a 2x2 table
    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(5), Inches(2))
    table = table_shape.table

    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "{{ key }}"
    table.cell(1, 1).text = "{{ value }}"

    path = os.path.join(tmp_dir, "table.pptx")
    prs.save(path)
    return path
