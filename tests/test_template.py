"""Tests for PptxTemplate core functionality."""

import os

import pytest
from pptx import Presentation

from pptxtpl import PptxTemplate
from pptxtpl.exceptions import InvalidTemplateError


class TestPptxTemplateInit:
    def test_load_valid_template(self, simple_template):
        tpl = PptxTemplate(simple_template)
        assert tpl.slides is not None

    def test_load_invalid_path_raises(self):
        with pytest.raises(InvalidTemplateError):
            PptxTemplate("/nonexistent/path.pptx")


class TestRenderSimple:
    def test_simple_variable(self, simple_template, tmp_dir):
        tpl = PptxTemplate(simple_template)
        tpl.render({"name": "World"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "World" in full_text
        assert "{{" not in full_text

    def test_multiple_variables(self, multi_var_template, tmp_dir):
        tpl = PptxTemplate(multi_var_template)
        tpl.render({"greeting": "Hi", "name": "Alice"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "Hi" in full_text
        assert "Alice" in full_text

    def test_empty_context(self, simple_template, tmp_dir):
        tpl = PptxTemplate(simple_template)
        tpl.render({"name": ""})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "Hello !" in full_text
        assert "{{" not in full_text


class TestRenderConditional:
    def test_condition_true(self, conditional_template, tmp_dir):
        tpl = PptxTemplate(conditional_template)
        tpl.render({"show_greeting": True, "name": "Bob"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "Hello Bob!" in full_text

    def test_condition_false(self, conditional_template, tmp_dir):
        tpl = PptxTemplate(conditional_template)
        tpl.render({"show_greeting": False, "name": "Bob"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "Hello" not in full_text


class TestRenderLoop:
    def test_for_loop(self, loop_template, tmp_dir):
        tpl = PptxTemplate(loop_template)
        tpl.render({"items": ["A", "B", "C"]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "A" in full_text
        assert "B" in full_text
        assert "C" in full_text

    def test_empty_loop(self, loop_template, tmp_dir):
        tpl = PptxTemplate(loop_template)
        tpl.render({"items": []})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "{%" not in full_text


class TestMultiSlide:
    def test_renders_all_slides(self, multi_slide_template, tmp_dir):
        tpl = PptxTemplate(multi_slide_template)
        tpl.render({"title": "Main Title", "subtitle": "Sub Title"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        texts_per_slide = []
        for slide in prs.slides:
            texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
            texts_per_slide.append(" ".join(texts))

        assert "Main Title" in texts_per_slide[0]
        assert "Sub Title" in texts_per_slide[1]


class TestTableTemplate:
    def test_table_variables(self, table_template, tmp_dir):
        tpl = PptxTemplate(table_template)
        tpl.render({"key": "Temperature", "value": "72°F"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                assert table.cell(1, 0).text == "Temperature"
                assert table.cell(1, 1).text == "72°F"
                break
        else:
            pytest.fail("No table found in output")


class TestGetUndeclaredVariables:
    def test_finds_variables(self, simple_template):
        tpl = PptxTemplate(simple_template)
        variables = tpl.get_undeclared_template_variables()
        assert "name" in variables

    def test_finds_multiple_variables(self, multi_var_template):
        tpl = PptxTemplate(multi_var_template)
        variables = tpl.get_undeclared_template_variables()
        assert "greeting" in variables
        assert "name" in variables

    def test_finds_variables_across_slides(self, multi_slide_template):
        tpl = PptxTemplate(multi_slide_template)
        variables = tpl.get_undeclared_template_variables()
        assert "title" in variables
        assert "subtitle" in variables
