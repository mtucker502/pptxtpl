"""Tests for multi-slide loop expansion ({%slide for%} spanning multiple slides)."""

import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptxtpl import PptxTemplate


def _get_slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


class TestMultiSlideLoopBasic:
    def test_two_slide_group_with_three_items(self, multi_slide_loop_template, tmp_dir):
        """A 2-slide group looped over 3 items produces 6 cloned slides."""
        tpl = PptxTemplate(multi_slide_loop_template)
        tpl.render({
            "title": "Report",
            "projects": [
                {"name": "Alpha", "detail": "Alpha details"},
                {"name": "Beta", "detail": "Beta details"},
                {"name": "Gamma", "detail": "Gamma details"},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # 1 title + 3×2 cloned + 1 closing = 8 slides
        assert len(prs.slides) == 8

        assert "Report" in _get_slide_text(prs.slides[0])
        # Project 1
        assert "Alpha" in _get_slide_text(prs.slides[1])
        assert "Alpha details" in _get_slide_text(prs.slides[2])
        # Project 2
        assert "Beta" in _get_slide_text(prs.slides[3])
        assert "Beta details" in _get_slide_text(prs.slides[4])
        # Project 3
        assert "Gamma" in _get_slide_text(prs.slides[5])
        assert "Gamma details" in _get_slide_text(prs.slides[6])

        assert "The End" in _get_slide_text(prs.slides[7])

    def test_no_jinja_tags_remain(self, multi_slide_loop_template, tmp_dir):
        tpl = PptxTemplate(multi_slide_loop_template)
        tpl.render({
            "title": "T",
            "projects": [{"name": "X", "detail": "Y"}],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        for slide in prs.slides:
            text = _get_slide_text(slide)
            assert "{%" not in text
            assert "{{" not in text


class TestMultiSlideLoopEmpty:
    def test_empty_list_removes_group(self, multi_slide_loop_template, tmp_dir):
        tpl = PptxTemplate(multi_slide_loop_template)
        tpl.render({
            "title": "Report",
            "projects": [],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # 1 title + 0 cloned + 1 closing = 2 slides
        assert len(prs.slides) == 2
        assert "Report" in _get_slide_text(prs.slides[0])
        assert "The End" in _get_slide_text(prs.slides[1])


class TestMultiSlideLoopSingle:
    def test_single_item_produces_one_group(self, multi_slide_loop_template, tmp_dir):
        tpl = PptxTemplate(multi_slide_loop_template)
        tpl.render({
            "title": "Report",
            "projects": [{"name": "Solo", "detail": "Solo detail"}],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        # 1 title + 1×2 cloned + 1 closing = 4 slides
        assert len(prs.slides) == 4
        assert "Solo" in _get_slide_text(prs.slides[1])
        assert "Solo detail" in _get_slide_text(prs.slides[2])


class TestMultiSlideLoopContext:
    def test_loop_helper_on_all_group_slides(self, tmp_dir):
        """The loop variable is available on every slide in the group."""
        prs = Presentation()

        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for item in items %}"
            "Slide A: {{ item.name }} {{ loop.index }}/{{ loop.length }}"
        )

        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "Slide B: {{ item.val }} {{ loop.index }}/{{ loop.length }}"
            "{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "ctx.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({
            "items": [
                {"name": "X", "val": "x1"},
                {"name": "Y", "val": "y1"},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 4

        assert "X" in _get_slide_text(result.slides[0])
        assert "1/2" in _get_slide_text(result.slides[0])
        assert "x1" in _get_slide_text(result.slides[1])
        assert "1/2" in _get_slide_text(result.slides[1])

        assert "Y" in _get_slide_text(result.slides[2])
        assert "2/2" in _get_slide_text(result.slides[2])
        assert "y1" in _get_slide_text(result.slides[3])
        assert "2/2" in _get_slide_text(result.slides[3])

    def test_loop_first_last(self, tmp_dir):
        prs = Presentation()

        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = (
            "{%slide for x in items %}"
            "{% if loop.first %}FIRST{% endif %}"
        )

        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = (
            "{% if loop.last %}LAST{% endif %}"
            "{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "fl.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({"items": [1, 2, 3]})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 6
        # Item 1: slides 0,1
        assert "FIRST" in _get_slide_text(result.slides[0])
        assert "LAST" not in _get_slide_text(result.slides[1])
        # Item 2: slides 2,3
        assert "FIRST" not in _get_slide_text(result.slides[2])
        assert "LAST" not in _get_slide_text(result.slides[3])
        # Item 3: slides 4,5
        assert "FIRST" not in _get_slide_text(result.slides[4])
        assert "LAST" in _get_slide_text(result.slides[5])


class TestMultiSlideLoopThreeSlides:
    def test_three_slide_group(self, tmp_dir):
        """A 3-slide group looped over 2 items produces 6 cloned slides."""
        prs = Presentation()

        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for p in projects %}Title: {{ p.name }}"
        )

        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "Body: {{ p.body }}"
        )

        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        s3.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "Footer: {{ p.footer }}{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "three_slide.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({
            "projects": [
                {"name": "A", "body": "Body A", "footer": "Foot A"},
                {"name": "B", "body": "Body B", "footer": "Foot B"},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 6

        assert "Title: A" in _get_slide_text(result.slides[0])
        assert "Body: Body A" in _get_slide_text(result.slides[1])
        assert "Footer: Foot A" in _get_slide_text(result.slides[2])
        assert "Title: B" in _get_slide_text(result.slides[3])
        assert "Body: Body B" in _get_slide_text(result.slides[4])
        assert "Footer: Foot B" in _get_slide_text(result.slides[5])


class TestMixedSingleAndMultiSlideLoops:
    def test_single_and_multi_slide_loops_together(self, tmp_dir):
        """A single-slide loop followed by a multi-slide loop."""
        prs = Presentation()

        # Static title
        s0 = prs.slides.add_slide(prs.slide_layouts[6])
        s0.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "Title"

        # Single-slide loop
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for a in alphas %}{{ a }}{%slide endfor %}"
        )

        # Multi-slide loop (2 slides)
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for p in projects %}Name: {{ p.name }}"
        )
        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        s3.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "Info: {{ p.info }}{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "mixed.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({
            "alphas": ["X", "Y"],
            "projects": [
                {"name": "P1", "info": "Info1"},
                {"name": "P2", "info": "Info2"},
                {"name": "P3", "info": "Info3"},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        # 1 title + 2 alpha + 3×2 project = 9
        assert len(result.slides) == 9
        assert "Title" in _get_slide_text(result.slides[0])
        assert "X" in _get_slide_text(result.slides[1])
        assert "Y" in _get_slide_text(result.slides[2])
        assert "P1" in _get_slide_text(result.slides[3])
        assert "Info1" in _get_slide_text(result.slides[4])
        assert "P2" in _get_slide_text(result.slides[5])
        assert "Info2" in _get_slide_text(result.slides[6])
        assert "P3" in _get_slide_text(result.slides[7])
        assert "Info3" in _get_slide_text(result.slides[8])


class TestMultiSlideLoopSpecialChars:
    def test_xml_special_chars(self, tmp_dir):
        prs = Presentation()

        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{%slide for item in items %}{{ item.name }}"
        )
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
            "{{ item.val }}{%slide endfor %}"
        )

        path = os.path.join(tmp_dir, "special.pptx")
        prs.save(path)

        tpl = PptxTemplate(path)
        tpl.render({
            "items": [
                {"name": "O'Brien & Sons", "val": "A < B"},
                {"name": "X > Y", "val": 'She said "hi"'},
            ],
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        result = Presentation(output)
        assert len(result.slides) == 4
        assert "O'Brien & Sons" in _get_slide_text(result.slides[0])
        assert "A < B" in _get_slide_text(result.slides[1])
        assert "X > Y" in _get_slide_text(result.slides[2])
        assert 'She said "hi"' in _get_slide_text(result.slides[3])
