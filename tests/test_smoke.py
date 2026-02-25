"""End-to-end smoke test for pptxtpl.

Generates a PowerPoint template and a JSON context file, renders
the template using pptxtpl, and verifies the output contains the
expected rendered content.
"""

import json
import os
import tempfile

from pptx import Presentation
from pptx.util import Inches

from pptxtpl import PptxTemplate, RichText


def test_smoke_end_to_end():
    """Full pipeline: create template + JSON context → render → verify output."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        # --- Step 1: Create a PowerPoint template with various Jinja2 features ---
        template_path = os.path.join(tmp_dir, "template.pptx")
        prs = Presentation()

        # Slide 1: Simple variable substitution
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        tb1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb1.text_frame.text = "{{ title }}"
        tb1b = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tb1b.text_frame.text = "By {{ author }}"

        # Slide 2: Conditional content
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tb2.text_frame.text = "{% if show_summary %}Summary: {{ summary }}{% endif %}"

        # Slide 3: Loop over items
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        tb3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
        tb3.text_frame.text = "{% for item in items %}{{ item.name }}: {{ item.desc }} {% endfor %}"

        # Slide 4: Table with variables
        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        tbl = slide4.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
        tbl.cell(0, 0).text = "Metric"
        tbl.cell(0, 1).text = "Value"
        tbl.cell(1, 0).text = "{{ metric1_name }}"
        tbl.cell(1, 1).text = "{{ metric1_value }}"
        tbl.cell(2, 0).text = "{{ metric2_name }}"
        tbl.cell(2, 1).text = "{{ metric2_value }}"

        prs.save(template_path)

        # --- Step 2: Create a JSON context file ---
        context_path = os.path.join(tmp_dir, "context.json")
        context_data = {
            "title": "Q4 2025 Report",
            "author": "Jane Smith",
            "show_summary": True,
            "summary": "Revenue up 15% year-over-year",
            "items": [
                {"name": "Product A", "desc": "Flagship product"},
                {"name": "Product B", "desc": "New launch"},
                {"name": "Product C", "desc": "Legacy support"},
            ],
            "metric1_name": "Revenue",
            "metric1_value": "$1.2M",
            "metric2_name": "Growth",
            "metric2_value": "15%",
        }
        with open(context_path, "w") as f:
            json.dump(context_data, f)

        # --- Step 3: Load JSON and render template ---
        with open(context_path) as f:
            context = json.load(f)

        tpl = PptxTemplate(template_path)

        # Verify undeclared variables are detected
        variables = tpl.get_undeclared_template_variables()
        assert "title" in variables
        assert "author" in variables
        assert "items" in variables

        tpl.render(context)
        output_path = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output_path)

        # --- Step 4: Open output and verify rendered content ---
        result = Presentation(output_path)
        assert len(result.slides) == 4

        # Slide 1: title and author
        slide1_texts = _get_slide_text(result.slides[0])
        assert "Q4 2025 Report" in slide1_texts
        assert "Jane Smith" in slide1_texts
        assert "{{" not in slide1_texts

        # Slide 2: conditional (should be shown since show_summary=True)
        slide2_texts = _get_slide_text(result.slides[1])
        assert "Revenue up 15% year-over-year" in slide2_texts
        assert "{%" not in slide2_texts

        # Slide 3: loop items
        slide3_texts = _get_slide_text(result.slides[2])
        assert "Product A" in slide3_texts
        assert "Flagship product" in slide3_texts
        assert "Product B" in slide3_texts
        assert "New launch" in slide3_texts
        assert "Product C" in slide3_texts
        assert "{%" not in slide3_texts

        # Slide 4: table
        slide4_table = None
        for shape in result.slides[3].shapes:
            if shape.has_table:
                slide4_table = shape.table
                break
        assert slide4_table is not None
        assert slide4_table.cell(1, 0).text == "Revenue"
        assert slide4_table.cell(1, 1).text == "$1.2M"
        assert slide4_table.cell(2, 0).text == "Growth"
        assert slide4_table.cell(2, 1).text == "15%"


def test_smoke_conditional_false():
    """Verify conditional content is omitted when condition is false."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        template_path = os.path.join(tmp_dir, "template.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tb.text_frame.text = "{% if show_summary %}Summary: {{ summary }}{% endif %}"
        prs.save(template_path)

        context = {"show_summary": False, "summary": "Should not appear"}
        tpl = PptxTemplate(template_path)
        tpl.render(context)

        output_path = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output_path)

        result = Presentation(output_path)
        text = _get_slide_text(result.slides[0])
        assert "Should not appear" not in text
        assert "Summary" not in text


def test_smoke_richtext_styled():
    """Verify RichText produces styled output in the rendered .pptx."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        template_path = os.path.join(tmp_dir, "template.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb.text_frame.text = "{{ styled_text }}"
        prs.save(template_path)

        rt = RichText("Important", bold=True, color="FF0000")
        rt.add(" note", italic=True)

        tpl = PptxTemplate(template_path)
        tpl.render({"styled_text": rt})
        output_path = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output_path)

        # Verify it's a valid pptx and has content
        result = Presentation(output_path)
        assert len(result.slides) == 1
        # The text should contain our content (RichText renders as XML runs)
        from lxml import etree
        slide_xml = etree.tostring(result.slides[0]._element, encoding="unicode")
        assert "Important" in slide_xml
        assert "note" in slide_xml


def _get_slide_text(slide) -> str:
    """Extract all text from a slide as a single string."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return " ".join(texts)
