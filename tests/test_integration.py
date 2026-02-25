"""Integration tests — end-to-end rendering with fragmented and complex templates."""

import os

import pytest
from pptx import Presentation

from pptxtpl import PptxTemplate, RichText, Listing


class TestFragmentedRuns:
    """Test that Jinja tags split across multiple XML runs still render correctly."""

    def test_fragmented_variable(self, fragmented_template, tmp_dir):
        tpl = PptxTemplate(fragmented_template)
        tpl.render({"name": "World"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "World" in full_text
        assert "{{" not in full_text


class TestRichTextRendering:
    def test_richtext_in_context(self, simple_template, tmp_dir):
        tpl = PptxTemplate(simple_template)
        rt = RichText("Bold Name", bold=True)
        tpl.render({"name": rt})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        # Verify the output is a valid pptx
        prs = Presentation(output)
        assert len(prs.slides) == 1


class TestListingRendering:
    def test_listing_in_context(self, simple_template, tmp_dir):
        tpl = PptxTemplate(simple_template)
        lst = Listing("Line 1\nLine 2")
        tpl.render({"name": lst})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 1


class TestComplexTemplates:
    def test_nested_conditionals_and_loops(self, tmp_dir):
        """Test a template with nested Jinja2 logic."""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = txBox.text_frame
        tf.text = "{% for item in items %}{% if item.active %}{{ item.name }}: {{ item.value }} {% endif %}{% endfor %}"
        tpl_path = os.path.join(tmp_dir, "complex.pptx")
        prs.save(tpl_path)

        tpl = PptxTemplate(tpl_path)
        tpl.render({
            "items": [
                {"name": "A", "value": 1, "active": True},
                {"name": "B", "value": 2, "active": False},
                {"name": "C", "value": 3, "active": True},
            ]
        })
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "A: 1" in full_text
        assert "B" not in full_text
        assert "C: 3" in full_text

    def test_special_characters_in_values(self, simple_template, tmp_dir):
        """Test that special XML characters in values are handled correctly."""
        tpl = PptxTemplate(simple_template)
        tpl.render({"name": "O'Brien & Sons <Ltd>"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "O'Brien & Sons <Ltd>" in full_text

    def test_numeric_values(self, simple_template, tmp_dir):
        """Test that numeric values render correctly."""
        tpl = PptxTemplate(simple_template)
        tpl.render({"name": 42})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full_text = " ".join(texts)
        assert "42" in full_text

    def test_render_then_save_preserves_slide_count(self, multi_slide_template, tmp_dir):
        """Verify slide count is preserved after rendering."""
        tpl = PptxTemplate(multi_slide_template)
        tpl.render({"title": "T", "subtitle": "S"})
        output = os.path.join(tmp_dir, "output.pptx")
        tpl.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2
