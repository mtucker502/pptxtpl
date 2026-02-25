"""Slide-level operations: cloning and deletion."""

import copy

from lxml import etree
from pptx.oxml.ns import qn


def clone_slide(prs, source_slide):
    """Clone a slide including its XML content and relationships.

    The clone is appended at the end of the presentation.
    Returns the new slide object.
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Collect source relationships (excluding the slide layout)
    src_rels = {}
    src_layout_rid = None
    for rel in source_slide.part.rels.values():
        if "slideLayout" in rel.reltype:
            src_layout_rid = rel.rId
        else:
            src_rels[rel.rId] = rel

    # Find the layout rId assigned to the new slide
    dst_layout_rid = None
    for rel in new_slide.part.rels.values():
        if "slideLayout" in rel.reltype:
            dst_layout_rid = rel.rId

    # Deep copy the source slide's XML element tree into the destination
    dst = new_slide._element
    for child in list(dst):
        dst.remove(child)
    for key in list(dst.attrib.keys()):
        del dst.attrib[key]

    for key, val in source_slide._element.attrib.items():
        dst.set(key, val)
    for child in source_slide._element:
        dst.append(copy.deepcopy(child))

    # Build rId remap table
    remap = {}
    if (
        src_layout_rid
        and dst_layout_rid
        and src_layout_rid != dst_layout_rid
    ):
        remap[src_layout_rid] = dst_layout_rid

    # Copy non-layout relationships and track rId changes
    for src_rid, rel in src_rels.items():
        if rel.is_external:
            new_rid = new_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        if new_rid != src_rid:
            remap[src_rid] = new_rid

    # Apply rId remapping to the copied XML
    if remap:
        _remap_rids(dst, remap)

    return new_slide


def delete_slide(prs, slide_index):
    """Remove a slide from the presentation by its zero-based index."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[slide_index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _remap_rids(element, remap):
    """Replace rId references in an XML element tree.

    Uses string replacement on the serialized XML to update all
    attribute values containing old rIds.
    """
    xml_str = etree.tostring(element, encoding="unicode")
    for old_rid, new_rid in remap.items():
        xml_str = xml_str.replace(f'"{old_rid}"', f'"{new_rid}"')

    new_element = etree.fromstring(xml_str.encode("utf-8"))

    # Replace content in-place
    for child in list(element):
        element.remove(child)
    for key in list(element.attrib.keys()):
        del element.attrib[key]
    for key, val in new_element.attrib.items():
        element.set(key, val)
    for child in new_element:
        element.append(child)
