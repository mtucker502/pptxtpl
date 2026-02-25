"""PptxTemplate — core render engine for pptxtpl.

Provides the main API for loading a .pptx template, rendering it with a
Jinja2 context dict, and saving the result.
"""

import re
from xml.sax.saxutils import escape

from lxml import etree
from jinja2 import Environment, BaseLoader, TemplateSyntaxError, meta
from pptx import Presentation
from pptx.oxml.ns import qn

from pptxtpl.xml_utils import preprocess_xml
from pptxtpl.richtext import RichText, Listing
from pptxtpl.slide_ops import clone_slide, delete_slide
from pptxtpl.exceptions import TemplateRenderError, InvalidTemplateError


# Regex for Jinja tags used to discover undeclared variables
_JINJA_TAG_RE = re.compile(r"(\{\{.*?\}\}|\{%.*?%\}|\{#.*?#\})", re.DOTALL)

# Slide-level loop tags: {%slide for VAR in EXPR %} and {%slide endfor %}
_SLIDE_FOR_RE = re.compile(
    r"\{%-?\s*slide\s+for\s+(\w+(?:\s*,\s*\w+)*)\s+in\s+(.*?)\s*-?%\}",
    re.DOTALL,
)
_SLIDE_ENDFOR_RE = re.compile(r"\{%-?\s*slide\s+endfor\s*-?%\}")


def _strip_slide_tags(xml: str) -> str:
    """Remove {%slide for%} and {%slide endfor%} tags from XML."""
    xml = _SLIDE_FOR_RE.sub("", xml)
    xml = _SLIDE_ENDFOR_RE.sub("", xml)
    return xml


def _escape_value(value):
    """Recursively XML-escape string values in a context value."""
    if isinstance(value, str):
        return escape(value)
    if isinstance(value, dict):
        return {k: _escape_value(v) for k, v in value.items()}
    if isinstance(value, (list, tuple)):
        return type(value)(_escape_value(v) for v in value)
    if isinstance(value, (RichText, Listing)):
        return str(value)
    return value


class PptxTemplate:
    """Load a PowerPoint template, render with Jinja2, and save.

    Usage::

        tpl = PptxTemplate("template.pptx")
        tpl.render({"name": "World"})
        tpl.save("output.pptx")
    """

    def __init__(self, template_path: str):
        self._template_path = template_path
        try:
            self._prs = Presentation(template_path)
        except Exception as exc:
            raise InvalidTemplateError(f"Cannot load template: {exc}") from exc

    @property
    def slides(self):
        """Access the presentation's slides."""
        return self._prs.slides

    def render(self, context: dict | None = None, jinja_env: Environment | None = None) -> None:
        """Render all slides with the given context dict.

        Args:
            context: Template variables dict. RichText and Listing values are
                     automatically converted to their XML/text representations.
            jinja_env: Optional custom Jinja2 Environment. If not provided,
                       a default environment is created.
        """
        if context is None:
            context = {}

        if jinja_env is None:
            jinja_env = Environment(loader=BaseLoader(), autoescape=False)
            jinja_env.globals.update({"RichText": RichText, "Listing": Listing})

        # Convert context values for Jinja2:
        # - RichText/Listing → their XML string representation (already escaped)
        # - Plain strings → XML-escaped to prevent invalid XML after rendering
        render_context = {}
        for key, value in context.items():
            if isinstance(value, (RichText, Listing)):
                render_context[key] = str(value)
            elif isinstance(value, str):
                render_context[key] = escape(value)
            else:
                render_context[key] = value

        # Phase 1: Expand slide-level loops (clones slides, modifies slide list)
        slide_contexts = self._expand_slide_loops(render_context, jinja_env)

        # Phase 2: Render each slide
        for i, slide in enumerate(self._prs.slides):
            ctx = render_context.copy()
            if i in slide_contexts:
                ctx.update(slide_contexts[i])
            self._render_slide(slide, ctx, jinja_env)

    def _expand_slide_loops(self, context: dict, jinja_env: Environment) -> dict:
        """Detect {%slide for%} tags and expand template slides into clones.

        Returns a dict mapping slide index → per-slide context overrides
        (the loop variable and a ``loop`` helper with index/first/last/length).
        """
        slide_contexts: dict[int, dict] = {}

        # First pass: identify template slides
        expansions = []
        for i, slide in enumerate(self._prs.slides):
            xml_str = etree.tostring(slide._element, encoding="unicode")
            xml_str = preprocess_xml(xml_str)
            match = _SLIDE_FOR_RE.search(xml_str)
            if match:
                expansions.append((i, match.group(1).strip(), match.group(2).strip()))

        if not expansions:
            return slide_contexts

        # Map context by rId so indices stay correct across multiple expansions
        rid_context: dict[str, dict] = {}
        sldIdLst = self._prs.slides._sldIdLst

        # Keep template sldIds in sldIdLst during expansion so that
        # add_slide (which uses len(sldIdLst) for partnames) never
        # generates duplicate names.  Remove them all at the end.
        deferred_removals: list[tuple] = []  # (template_sldId, rId)

        # Process in reverse order so that earlier indices remain valid
        for slide_idx, var_names_str, iterable_expr in reversed(expansions):
            # Evaluate the iterable expression using Jinja2
            try:
                expr_fn = jinja_env.compile_expression(iterable_expr)
                items = list(expr_fn(**context))
            except Exception as exc:
                raise TemplateRenderError(
                    f"Cannot evaluate slide loop iterable '{iterable_expr}': {exc}"
                ) from exc

            template_sldId = sldIdLst[slide_idx]
            rId = template_sldId.get(qn("r:id"))

            if not items:
                deferred_removals.append((template_sldId, rId))
                continue

            source_slide = self._prs.slides[slide_idx]

            # Clone the slide for each item — do NOT touch sldIdLst between
            # clones, because add_slide uses len(prs.slides) to generate
            # unique part names.
            n_before = len(list(sldIdLst))
            for _ in items:
                clone_slide(self._prs, source_slide)

            # Collect the clone sldIds (they were appended at the end)
            n_clones = len(items)
            clone_sldIds = [sldIdLst[n_before + i] for i in range(n_clones)]

            # Remove clones from the end of sldIdLst
            for sldId in reversed(clone_sldIds):
                sldIdLst.remove(sldId)

            # Insert clones just before the template slide
            template_pos = list(sldIdLst).index(template_sldId)
            for i, clone_sldId in enumerate(clone_sldIds):
                sldIdLst.insert(template_pos + i, clone_sldId)

            # Mark template for deferred removal (keep in sldIdLst for now)
            deferred_removals.append((template_sldId, rId))

            # Store per-slide context keyed by rId
            var_list = [v.strip() for v in var_names_str.split(",")]
            n_items = len(items)
            for i, (clone_sldId, item) in enumerate(zip(clone_sldIds, items)):
                clone_rId = clone_sldId.get(qn("r:id"))
                ctx: dict = {}

                # Bind loop variable(s) — with recursive XML escaping
                escaped_item = _escape_value(item)
                if len(var_list) == 1:
                    ctx[var_list[0]] = escaped_item
                else:
                    for var_name, val in zip(var_list, escaped_item):
                        ctx[var_name] = val

                # Provide a loop helper (mirrors Jinja2's loop variable)
                ctx["loop"] = {
                    "index": i + 1,
                    "index0": i,
                    "first": i == 0,
                    "last": i == n_items - 1,
                    "length": n_items,
                }

                rid_context[clone_rId] = ctx

        # Now remove all template slides and drop their relationships
        for template_sldId, rId in deferred_removals:
            sldIdLst.remove(template_sldId)
            self._prs.part.drop_rel(rId)

        # Convert rId-keyed contexts to index-keyed using the final slide order
        for i, sldId in enumerate(sldIdLst):
            rId = sldId.get(qn("r:id"))
            if rId in rid_context:
                slide_contexts[i] = rid_context[rId]

        return slide_contexts

    def _render_slide(self, slide, context: dict, jinja_env: Environment) -> None:
        """Render a single slide's XML through the Jinja2 pipeline."""
        # Get the slide's XML element
        slide_element = slide._element

        # Serialize to XML string
        xml_str = etree.tostring(slide_element, encoding="unicode")

        # Preprocess: fix fragmented delimiters, strip internal tags, etc.
        xml_str = preprocess_xml(xml_str)

        # Strip slide-level loop tags (already handled by _expand_slide_loops)
        xml_str = _strip_slide_tags(xml_str)

        # Check if there are any Jinja tags after preprocessing
        if not _JINJA_TAG_RE.search(xml_str):
            return  # No templates on this slide

        # Render with Jinja2
        try:
            template = jinja_env.from_string(xml_str)
            rendered_xml = template.render(context)
        except TemplateSyntaxError as exc:
            raise TemplateRenderError(
                f"Jinja2 syntax error on slide: {exc}"
            ) from exc
        except Exception as exc:
            raise TemplateRenderError(
                f"Rendering failed on slide: {exc}"
            ) from exc

        # Post-process: convert \n to line breaks, \a to paragraph breaks
        rendered_xml = self._post_process(rendered_xml)

        # Parse the rendered XML back into an element tree
        try:
            new_element = etree.fromstring(rendered_xml.encode("utf-8"))
        except etree.XMLSyntaxError as exc:
            raise TemplateRenderError(
                f"Rendered XML is invalid: {exc}"
            ) from exc

        # Replace the slide's element tree
        parent = slide_element.getparent()
        if parent is not None:
            parent.replace(slide_element, new_element)
            slide._element = new_element
        else:
            # Slide is the root element — replace children
            slide_element.clear()
            for attr_name, attr_val in new_element.attrib.items():
                slide_element.set(attr_name, attr_val)
            for child in new_element:
                slide_element.append(child)

    def _post_process(self, xml: str) -> str:
        """Convert escape sequences in rendered text to PowerPoint XML.

        - ``\\n`` inside <a:t> elements becomes ``<a:br/>``
        - ``\\a`` becomes a paragraph break (closes and reopens <a:p>)
        """
        # Handle \n → line break within <a:t> elements
        # We replace \n in text content with </a:t></a:r><a:br/><a:r><a:t>
        xml = self._replace_newlines_in_text(xml)
        return xml

    def _replace_newlines_in_text(self, xml: str) -> str:
        """Replace literal \\n characters inside <a:t> elements with <a:br/> elements."""

        def _replace_in_at(match: re.Match) -> str:
            opening = match.group(1)
            content = match.group(2)
            if "\n" not in content:
                return match.group(0)
            # Split on \n and join with line break XML
            parts = content.split("\n")
            result = opening + parts[0] + "</a:t></a:r><a:br/><a:r>" + ("<a:t>".join(
                p + "</a:t></a:r><a:br/><a:r>" for p in parts[1:-1]
            )) + "<a:t>" + parts[-1] + "</a:t>"
            return result

        return re.sub(r"(<a:t[^>]*>)(.*?)</a:t>", _replace_in_at, xml, flags=re.DOTALL)

    def save(self, output_path: str) -> None:
        """Save the rendered presentation to a file."""
        self._prs.save(output_path)

    def get_undeclared_template_variables(
        self, jinja_env: Environment | None = None
    ) -> set[str]:
        """Find all undeclared variables across all slides.

        Returns a set of variable names that appear in Jinja2 expressions
        but are not defined in the environment's globals.
        """
        if jinja_env is None:
            jinja_env = Environment(loader=BaseLoader(), autoescape=False)

        all_vars: set[str] = set()

        for slide in self._prs.slides:
            xml_str = etree.tostring(slide._element, encoding="unicode")
            xml_str = preprocess_xml(xml_str)
            xml_str = _strip_slide_tags(xml_str)

            try:
                ast = jinja_env.parse(xml_str)
                variables = meta.find_undeclared_variables(ast)
                all_vars.update(variables)
            except TemplateSyntaxError:
                pass  # Skip slides with syntax errors

        return all_vars
