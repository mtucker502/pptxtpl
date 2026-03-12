"""Render a multi-slide loop: two slides per project.

Demonstrates {%slide for%} on one slide and {%slide endfor%} on a later
slide so that each loop iteration clones the entire group.

Usage:
    python examples/multi_slide_loop_render.py
"""

import os

from pptx import Presentation

from pptxtpl import PptxTemplate

here = os.path.dirname(__file__)
template_path = os.path.join(here, "multi_slide_loop_template.pptx")
output_path = os.path.join(here, "multi_slide_loop_output.pptx")

# -- Context ------------------------------------------------------------------
context = {
    "report": {
        "title": "Q1 2026 Project Portfolio",
        "date": "2026-03-12",
        "closing": "End of Report",
    },
    "projects": [
        {
            "name": "Atlas",
            "status": "On track",
            "lead": "Alice Chen",
            "summary": "Backend API platform for internal services.",
            "description": (
                "Atlas provides a unified REST/gRPC gateway for all internal "
                "microservices. The Q1 milestone focuses on rate limiting, "
                "circuit breakers, and observability dashboards."
            ),
            "tags": ["backend", "infrastructure", "Q1"],
        },
        {
            "name": "Beacon",
            "status": "At risk",
            "lead": "Bob Martinez",
            "summary": "Customer-facing notification system overhaul.",
            "description": (
                "Beacon replaces the legacy email/SMS pipeline with a unified "
                "notification service supporting push, email, SMS, and in-app "
                "channels. Delayed by vendor contract negotiations."
            ),
            "tags": ["frontend", "notifications", "Q1"],
        },
        {
            "name": "Comet",
            "status": "Complete",
            "lead": "Carol Nguyen",
            "summary": "Data pipeline migration from Airflow to Dagster.",
            "description": (
                "Comet migrated 47 DAGs from Airflow 1.x to Dagster, "
                "reducing pipeline failures by 62% and cutting average "
                "orchestration latency from 12s to 3s."
            ),
            "tags": ["data", "infrastructure", "Q4-carryover"],
        },
    ],
}

print(f"Projects to render: {len(context['projects'])}")

# -- Render --------------------------------------------------------------------
tpl = PptxTemplate(template_path)
tpl.render(context)
tpl.save(output_path)

print(f"Output saved to {output_path}")

# -- Verify --------------------------------------------------------------------
result = Presentation(output_path)
print(f"Total slides: {len(result.slides)}")
print(f"  Expected: 1 title + {len(context['projects'])}x2 project + 1 closing "
      f"= {1 + len(context['projects']) * 2 + 1}")

for i, slide in enumerate(result.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    preview = " | ".join(t for t in texts if t.strip())
    if len(preview) > 120:
        preview = preview[:120] + "..."
    print(f"  Slide {i}: {preview}")
