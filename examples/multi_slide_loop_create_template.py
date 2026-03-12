"""Create a template demonstrating multi-slide loops.

Each project gets TWO slides per iteration:
  1. Project Summary — name, status, and lead
  2. Project Details — description, milestones, and tags

Slides:
  1. Title (static)
  2. Project Summary  (cloned as a group, once per project)
  3. Project Details   ↑
  4. Closing (static)

Run this once to produce examples/multi_slide_loop_template.pptx.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

here = os.path.dirname(__file__)
output_path = os.path.join(here, "multi_slide_loop_template.pptx")

prs = Presentation()

NAVY = RGBColor(0x1A, 0x3C, 0x6E)
ACCENT = RGBColor(0xE8, 0x8D, 0x2A)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF2, 0xF4, 0xF7)


# -- Slide 1: Title -----------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

tb = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "{{ report.title }}"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

tb2 = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.5))
p = tb2.text_frame.paragraphs[0]
p.text = "{{ report.date }}"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER


# -- Slide 2: Project Summary (loop start) ------------------------------------
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Loop-open directive (hidden)
tag_open = slide2.shapes.add_textbox(Inches(0), Inches(0), Inches(0.01), Inches(0.01))
tag_open.text_frame.text = "{%slide for project in projects %}"

# Slide counter (top-right)
counter = slide2.shapes.add_textbox(Inches(7), Inches(0.2), Inches(2.5), Inches(0.3))
p = counter.text_frame.paragraphs[0]
p.text = "Project {{ loop.index }} of {{ loop.length }}"
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.RIGHT

# "SUMMARY" label
lbl = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.3))
p = lbl.text_frame.paragraphs[0]
p.text = "SUMMARY"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = ACCENT

# Project name badge
badge = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(0.8), Inches(4), Inches(0.6),
)
badge.fill.solid()
badge.fill.fore_color.rgb = NAVY
badge.line.fill.background()
p = badge.text_frame.paragraphs[0]
p.text = "{{ project.name }}"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT
badge.text_frame.margin_left = Inches(0.15)

# Status
status_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
tf = status_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Status: {{ project.status }}"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

# Lead
lead_box = slide2.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(0.5))
tf = lead_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Lead: {{ project.lead }}"
p.font.size = Pt(14)
p.font.color.rgb = GRAY

# Accent bar
bar = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.1), Inches(9), Emu(36000),
)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Summary text
summary_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(2))
tf = summary_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "{{ project.summary }}"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# -- Slide 3: Project Details (loop end) --------------------------------------
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Slide counter (top-right)
counter2 = slide3.shapes.add_textbox(Inches(7), Inches(0.2), Inches(2.5), Inches(0.3))
p = counter2.text_frame.paragraphs[0]
p.text = "Project {{ loop.index }} of {{ loop.length }}"
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.RIGHT

# "DETAILS" label
lbl2 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.3))
p = lbl2.text_frame.paragraphs[0]
p.text = "DETAILS"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = ACCENT

# Project name (repeated for context)
name_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.6))
p = name_box.text_frame.paragraphs[0]
p.text = "{{ project.name }}"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = NAVY

# Accent bar
bar2 = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Emu(36000),
)
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT
bar2.line.fill.background()

# Description label + text
desc_lbl = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(2), Inches(0.3))
p = desc_lbl.text_frame.paragraphs[0]
p.text = "Description"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = GRAY

desc_box = slide3.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "{{ project.description }}"
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Tags label + list
tags_lbl = slide3.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(2), Inches(0.3))
p = tags_lbl.text_frame.paragraphs[0]
p.text = "Tags"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = GRAY

tags_box = slide3.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.5))
tf = tags_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '{{ project.tags | join(", ") }}'
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Loop-close directive (hidden)
tag_close = slide3.shapes.add_textbox(Inches(0), Inches(0), Inches(0.01), Inches(0.01))
tag_close.text_frame.text = "{%slide endfor %}"


# -- Slide 4: Closing ---------------------------------------------------------
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

end_box = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = end_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "{{ report.closing }}"
p.font.size = Pt(24)
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER


prs.save(output_path)
print(f"Template saved to {output_path}")
